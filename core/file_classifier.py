"""大嘴怪 — 文件类型识别引擎

基于 python-magic MIME 识别 + 各类型元数据提取 + 文本内容提取。
输出统一 FileInfo 数据类，边界情况下不崩溃。
"""
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional
import json
import hashlib
import mimetypes

from core.logger import get_logger

logger = get_logger()

# ── 扩展名 → MIME 类型补充映射（magic 失败时回退）───
EXT_TO_MIME = {
    # 文档
    ".pdf": "application/pdf",
    ".doc": "application/msword",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".xls": "application/vnd.ms-excel",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".ppt": "application/vnd.ms-powerpoint",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".txt": "text/plain",
    ".md": "text/markdown",
    ".rst": "text/x-rst",
    ".csv": "text/csv",
    ".json": "application/json",
    ".xml": "application/xml",
    ".yaml": "text/yaml",
    ".yml": "text/yaml",
    ".toml": "text/toml",
    ".py": "text/x-python",
    ".js": "text/javascript",
    ".ts": "text/typescript",
    ".java": "text/x-java",
    ".go": "text/x-go",
    ".rs": "text/x-rust",
    ".cpp": "text/x-c++",
    ".c": "text/x-c",
    ".h": "text/x-c",
    ".html": "text/html",
    ".css": "text/css",
    ".sql": "text/x-sql",
    ".sh": "text/x-shellscript",
    # 图片
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".gif": "image/gif",
    ".bmp": "image/bmp",
    ".webp": "image/webp",
    ".svg": "image/svg+xml",
    ".ico": "image/x-icon",
    ".tiff": "image/tiff",
    ".tif": "image/tiff",
    ".psd": "image/vnd.adobe.photoshop",
    # 音视频
    ".mp3": "audio/mpeg",
    ".wav": "audio/wav",
    ".flac": "audio/flac",
    ".aac": "audio/aac",
    ".ogg": "audio/ogg",
    ".wma": "audio/x-ms-wma",
    ".mp4": "video/mp4",
    ".avi": "video/x-msvideo",
    ".mkv": "video/x-matroska",
    ".mov": "video/quicktime",
    ".wmv": "video/x-ms-wmv",
    ".flv": "video/x-flv",
    ".webm": "video/webm",
    ".m4a": "audio/mp4",
    ".m4v": "video/mp4",
    # 压缩包
    ".zip": "application/zip",
    ".rar": "application/x-rar-compressed",
    ".7z": "application/x-7z-compressed",
    ".tar": "application/x-tar",
    ".gz": "application/gzip",
    ".bz2": "application/x-bzip2",
    ".xz": "application/x-xz",
}

# ── MIME 类型 → 大类映射 ────────────────────────────
MIME_TO_TYPE = {
    # 文档
    "application/pdf": "document",
    "application/msword": "document",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "document",
    "application/vnd.ms-excel": "document",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "document",
    "application/vnd.ms-powerpoint": "document",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "document",
    "application/json": "document",
    "application/xml": "document",
    # 文本
    "text/plain": "document",
    "text/markdown": "document",
    "text/csv": "document",
    "text/html": "document",
    "text/x-python": "document",
    "text/javascript": "document",
    "text/typescript": "document",
    "text/x-java": "document",
    "text/x-go": "document",
    "text/x-rust": "document",
    "text/x-c++": "document",
    "text/x-c": "document",
    "text/css": "document",
    "text/x-sql": "document",
    "text/x-shellscript": "document",
    "text/yaml": "document",
    "text/toml": "document",
    # 图片
    "image/png": "image",
    "image/jpeg": "image",
    "image/gif": "image",
    "image/bmp": "image",
    "image/webp": "image",
    "image/svg+xml": "image",
    "image/x-icon": "image",
    "image/tiff": "image",
    # 音频
    "audio/mpeg": "audio",
    "audio/wav": "audio",
    "audio/flac": "audio",
    "audio/aac": "audio",
    "audio/ogg": "audio",
    "audio/mp4": "audio",
    # 视频
    "video/mp4": "video",
    "video/x-msvideo": "video",
    "video/x-matroska": "video",
    "video/quicktime": "video",
    "video/webm": "video",
    # 压缩包
    "application/zip": "archive",
    "application/x-rar-compressed": "archive",
    "application/x-7z-compressed": "archive",
    "application/x-tar": "archive",
    "application/gzip": "archive",
    "application/x-bzip2": "archive",
    "application/x-xz": "archive",
}

# ── FileInfo 数据类 ───────────────────────────────

@dataclass
class FileInfo:
    """统一的文件信息数据类"""
    file_path: str = ""
    file_name: str = ""
    file_size: int = 0
    checksum: str = ""                        # MD5
    mime_type: str = ""
    content_type: str = "other"              # document/image/audio/video/archive/other
    metadata: dict = field(default_factory=dict)
    text_preview: str = ""                   # 提取的文本内容
    text_length: int = 0
    error: str = ""                          # 处理错误信息（空 = 正常）

    @property
    def title(self) -> str:
        """文件名去掉扩展名"""
        return Path(self.file_name).stem if self.file_name else ""

    def to_dict(self) -> dict:
        return {
            "file_path": self.file_path,
            "file_name": self.file_name,
            "file_size": self.file_size,
            "checksum": self.checksum,
            "mime_type": self.mime_type,
            "content_type": self.content_type,
            "metadata": json.dumps(self.metadata, ensure_ascii=False),
            "text_preview": self.text_preview,
            "text_length": self.text_length,
            "error": self.error,
        }


# ── FileClassifier ────────────────────────────────

class FileClassifier:
    """文件类型识别引擎"""

    def identify(self, file_path: str) -> FileInfo:
        """主入口：识别任意文件，返回 FileInfo"""
        path = Path(file_path)
        info = FileInfo(
            file_path=str(path),
            file_name=path.name,
            file_size=path.stat().st_size if path.exists() else 0,
        )

        if not path.exists():
            info.error = f"文件不存在: {file_path}"
            return info

        try:
            # 1. 计算 MD5
            info.checksum = self._compute_md5(path)

            # 2. 识别 MIME 类型
            info.mime_type = self._detect_mime(path)

            # 3. MIME → content_type
            info.content_type = self._mime_to_type(info.mime_type)

            # 4. 提取元数据 + 文本
            self._extract_by_type(path, info)

        except Exception as e:
            logger.warning(f"FileClassifier 处理异常: {file_path} — {e}")
            info.error = str(e)

        return info

    # ── MIME 检测 ─────────────────────────────────
    def _detect_mime(self, path: Path) -> str:
        """MIME 类型检测，带降级"""
        try:
            import magic
            mime = magic.from_file(str(path), mime=True)
            if mime and mime != "application/octet-stream":
                return mime
        except (ImportError, Exception):
            pass  # python-magic 不可用

        # 降级：扩展名 → MIME
        ext = path.suffix.lower()
        return EXT_TO_MIME.get(ext, "application/octet-stream")

    def _mime_to_type(self, mime: str) -> str:
        """MIME → 大类"""
        if mime in MIME_TO_TYPE:
            return MIME_TO_TYPE[mime]
        # 前缀匹配
        if mime.startswith("text/"):
            return "document"
        if mime.startswith("image/"):
            return "image"
        if mime.startswith("audio/"):
            return "audio"
        if mime.startswith("video/"):
            return "video"
        return "other"

    # ── 分类型提取 ────────────────────────────────
    def _extract_by_type(self, path: Path, info: FileInfo) -> None:
        """根据文件类型分派到对应的提取器"""
        ext = path.suffix.lower()
        ct = info.content_type
        mime = info.mime_type

        if ct == "document":
            self._extract_document(path, info, ext, mime)
        elif ct == "image":
            self._extract_image(path, info, ext)
        elif ct in ("audio", "video"):
            self._extract_media(path, info, ct)
        elif ct == "archive":
            self._extract_archive(path, info)
        else:
            # 未知类型：尝试当文本读
            self._try_read_as_text(path, info)

    # ── 文档提取 ──────────────────────────────────
    def _extract_document(self, path: Path, info: FileInfo, ext: str, mime: str) -> None:
        """文档元数据 + 文本提取"""
        reader = self._get_doc_reader(ext, mime)
        if reader:
            try:
                reader(path, info)
                return
            except Exception as e:
                logger.debug(f"文档读取器失败 ({reader.__name__}): {e}")
        # 降级：纯文本读取
        self._try_read_as_text(path, info)

    def _get_doc_reader(self, ext: str, mime: str):
        """获取文档读取函数"""
        if ext == ".pdf":
            return self._read_pdf
        if ext in (".docx",):
            return self._read_docx
        if ext in (".pptx",):
            return self._read_pptx
        if ext in (".xlsx", ".xls"):
            return self._read_xlsx
        if mime == "application/json":
            return self._read_json
        return None

    def _read_pdf(self, path: Path, info: FileInfo) -> None:
        """PDF 提取"""
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(str(path))
            info.metadata["pages"] = len(reader.pages)
            if reader.metadata:
                info.metadata["pdf_title"] = reader.metadata.get("/Title", "")
                info.metadata["author"] = reader.metadata.get("/Author", "")

            # 提取全文文本
            texts = []
            for page in reader.pages[:20]:  # 最多前 20 页
                t = page.extract_text()
                if t:
                    texts.append(t)
            full_text = "\n".join(texts)
            info.text_preview = full_text[:5000]
            info.text_length = len(full_text)
        except Exception:
            # PyPDF2 失败，尝试 pdfplumber
            try:
                import pdfplumber
                with pdfplumber.open(str(path)) as pdf:
                    info.metadata["pages"] = len(pdf.pages)
                    texts = []
                    for page in pdf.pages[:20]:
                        t = page.extract_text()
                        if t:
                            texts.append(t)
                    full_text = "\n".join(texts)
                    info.text_preview = full_text[:5000]
                    info.text_length = len(full_text)
            except Exception:
                self._try_read_as_text(path, info)

    def _read_docx(self, path: Path, info: FileInfo) -> None:
        """Word 文档提取"""
        from docx import Document
        doc = Document(str(path))
        info.metadata["paragraphs"] = len(doc.paragraphs)
        # 提取属性
        props = doc.core_properties
        info.metadata["author"] = props.author or ""
        info.metadata["doc_title"] = props.title or ""

        texts = [p.text for p in doc.paragraphs if p.text.strip()]
        full_text = "\n".join(texts)
        info.text_preview = full_text[:5000]
        info.text_length = len(full_text)

    def _read_pptx(self, path: Path, info: FileInfo) -> None:
        """PPT 提取"""
        from pptx import Presentation
        prs = Presentation(str(path))
        info.metadata["slides"] = len(prs.slides)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            texts.append(para.text)
        full_text = "\n".join(texts)
        info.text_preview = full_text[:5000]
        info.text_length = len(full_text)

    def _read_xlsx(self, path: Path, info: FileInfo) -> None:
        """Excel 提取"""
        from openpyxl import load_workbook
        wb = load_workbook(str(path), read_only=True, data_only=True)
        info.metadata["sheets"] = wb.sheetnames
        texts = []
        for sheet_name in wb.sheetnames[:5]:  # 最多前 5 个 sheet
            ws = wb[sheet_name]
            for row in ws.iter_rows(values_only=True):
                row_text = " ".join(str(c) for c in row if c is not None)
                if row_text.strip():
                    texts.append(row_text)
                if len(texts) > 500:
                    break
        full_text = "\n".join(texts)
        info.text_preview = full_text[:5000]
        info.text_length = len(full_text)
        wb.close()

    def _read_json(self, path: Path, info: FileInfo) -> None:
        """JSON 提取"""
        try:
            text = path.read_text(encoding="utf-8")
            data = json.loads(text)
            info.text_preview = json.dumps(data, ensure_ascii=False, indent=2)[:5000]
            info.text_length = len(text)
        except Exception:
            self._try_read_as_text(path, info)

    # ── 图片提取 ──────────────────────────────────
    def _extract_image(self, path: Path, info: FileInfo, ext: str) -> None:
        """图片元数据 + 可选 OCR"""
        try:
            from PIL import Image, ExifTags
            with Image.open(str(path)) as img:
                info.metadata["width"] = img.width
                info.metadata["height"] = img.height
                info.metadata["format"] = img.format or ""
                info.metadata["mode"] = img.mode

                # EXIF
                exif_data = img.getexif()
                if exif_data:
                    for tag_id, value in exif_data.items():
                        tag_name = ExifTags.TAGS.get(tag_id, str(tag_id))
                        if tag_name in ("DateTime", "Make", "Model", "GPSInfo", "ImageDescription"):
                            info.metadata[f"exif_{tag_name}"] = str(value)

                # OCR（可选）
                info.text_preview = self._ocr_image(path)
                info.text_length = len(info.text_preview)

        except ImportError:
            logger.warning("Pillow 不可用，跳过图片元数据提取")
        except Exception as e:
            info.metadata["error"] = str(e)

    def _ocr_image(self, path: Path) -> str:
        """OCR 文字识别（可选）"""
        try:
            import pytesseract
            from PIL import Image
            img = Image.open(str(path))
            text = pytesseract.image_to_string(img, lang="chi_sim+eng")
            return text.strip()[:2000]
        except ImportError:
            return ""
        except Exception:
            return ""

    # ── 音视频提取 ────────────────────────────────
    def _extract_media(self, path: Path, info: FileInfo, ct: str) -> None:
        """音视频元数据提取"""
        try:
            import subprocess
            import json as _json
            result = subprocess.run(
                ["ffprobe", "-v", "quiet", "-print_format", "json",
                 "-show_format", "-show_streams", str(path)],
                capture_output=True, text=True, timeout=30,
            )
            if result.returncode == 0 and result.stdout:
                data = _json.loads(result.stdout)
                fmt = data.get("format", {})
                info.metadata["duration"] = fmt.get("duration", "")
                info.metadata["bit_rate"] = fmt.get("bit_rate", "")
                info.metadata["format_name"] = fmt.get("format_name", "")

                for stream in data.get("streams", []):
                    codec_type = stream.get("codec_type", "")
                    if codec_type == "video":
                        info.metadata["video_codec"] = stream.get("codec_name", "")
                        info.metadata["video_width"] = stream.get("width", 0)
                        info.metadata["video_height"] = stream.get("height", 0)
                        info.metadata["fps"] = stream.get("r_frame_rate", "")
                    elif codec_type == "audio":
                        info.metadata["audio_codec"] = stream.get("codec_name", "")
                        info.metadata["sample_rate"] = stream.get("sample_rate", "")
                        info.metadata["channels"] = stream.get("channels", 0)

            info.text_preview = (
                f"[{ct}] {path.name}\n"
                f"时长: {info.metadata.get('duration', '未知')}\n"
                f"格式: {info.metadata.get('format_name', '未知')}"
            )
            info.text_length = len(info.text_preview)

        except FileNotFoundError:
            info.metadata["note"] = "ffprobe 不可用，仅记录基本文件信息"
            info.text_preview = f"[{ct}] {path.name}"
            info.text_length = len(info.text_preview)
            info.error = "ffprobe 未安装，无法提取音视频元数据"
        except Exception as e:
            info.metadata["error"] = str(e)
            info.text_preview = f"[{ct}] {path.name}"
            info.text_length = len(info.text_preview)

    # ── 压缩包 ─────────────────────────────────────
    def _extract_archive(self, path: Path, info: FileInfo) -> None:
        """压缩包元数据"""
        info.text_preview = f"[压缩包] {path.name}"
        info.text_length = len(info.text_preview)
        # 可扩展：列出包内文件清单

    # ── 通用文本读取 ──────────────────────────────
    def _try_read_as_text(self, path: Path, info: FileInfo) -> None:
        """尝试以文本方式读取文件"""
        # 仅对小于 10MB 的文件尝试
        if info.file_size > 10 * 1024 * 1024:
            info.text_preview = f"[大文件] {path.name}"
            info.text_length = len(info.text_preview)
            return

        encodings = ["utf-8", "gbk", "gb2312", "latin-1"]
        for enc in encodings:
            try:
                text = path.read_text(encoding=enc)
                info.text_preview = text[:5000]
                info.text_length = len(text)
                info.metadata["encoding"] = enc
                return
            except (UnicodeDecodeError, UnicodeError):
                continue

        info.text_preview = f"[二进制] {path.name}"
        info.text_length = len(info.text_preview)

    # ── MD5 ────────────────────────────────────────
    @staticmethod
    def _compute_md5(path: Path) -> str:
        h = hashlib.md5()
        with open(str(path), "rb") as f:
            for chunk in iter(lambda: f.read(65536), b""):
                h.update(chunk)
        return h.hexdigest()

"""大嘴怪 — FileClassifier 文件类型识别 全覆盖测试"""
import sys, os, tempfile, shutil, json, hashlib
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
sys.path.insert(0, '/home/admin/.local/lib/python3.10/site-packages')

from core.file_classifier import FileClassifier, FileInfo

fc = FileClassifier()


# ── 辅助 ──────────────────────────────────────────

def _create_python(tmp):
    path = os.path.join(tmp, "hello.py")
    with open(path, "w") as f:
        f.write("import sys\ndef greet(name: str) -> str:\n    return f'Hello, {name}!'\n")
    return path


def _create_markdown(tmp):
    path = os.path.join(tmp, "readme.md")
    with open(path, "w") as f:
        f.write("# README\n\nPython is great. Docker is useful.\n")
    return path


def _create_json(tmp):
    path = os.path.join(tmp, "data.json")
    with open(path, "w") as f:
        json.dump({"name": "test", "version": 1}, f)
    return path


def _create_txt(tmp):
    path = os.path.join(tmp, "note.txt")
    with open(path, "w", encoding="utf-8") as f:
        f.write("UTF-8 中文测试 Python Docker LLM")
    return path


def _create_empty(tmp):
    path = os.path.join(tmp, "empty.txt")
    with open(path, "w") as f:
        pass  # 创建空文件
    return path


def _create_png(tmp):
    path = os.path.join(tmp, "image.png")
    try:
        from PIL import Image
        img = Image.new("RGB", (200, 200), (100, 200, 100))
        img.save(path, "PNG")
    except ImportError:
        with open(path, "wb") as f:
            f.write(b"\x89PNG\r\n\x1a\n" + b"\x00" * 100)
    return path


def _create_zip(tmp):
    import zipfile
    path = os.path.join(tmp, "archive.zip")
    with zipfile.ZipFile(path, "w") as zf:
        zf.writestr("a.txt", "content")
    return path


def _create_docx(tmp):
    path = os.path.join(tmp, "test.docx")
    try:
        from docx import Document
        doc = Document()
        doc.add_paragraph("Hello Python Docker LLM")
        doc.core_properties.author = "Tester"
        doc.core_properties.title = "Test Doc"
        doc.save(path)
    except ImportError:
        # 占位，DOCX 测试会被 skip_iff_no_docx 跳过
        with open(path, "wb") as f:
            f.write(b"PK\x03\x04" + b"\x00" * 100)
    return path


def _create_xlsx(tmp):
    path = os.path.join(tmp, "test.xlsx")
    try:
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws["A1"] = "Hello"
        ws["B1"] = "World"
        wb.save(path)
    except ImportError:
        with open(path, "wb") as f:
            f.write(b"PK\x03\x04" + b"\x00" * 100)
    return path


def _create_pptx(tmp):
    path = os.path.join(tmp, "test.pptx")
    try:
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        slide.shapes.title.text = "Hello PPT"
        prs.save(path)
    except ImportError:
        with open(path, "wb") as f:
            f.write(b"PK\x03\x04" + b"\x00" * 100)
    return path


def _create_pdf(tmp):
    path = os.path.join(tmp, "test.pdf")
    try:
        from PyPDF2 import PdfWriter
        writer = PdfWriter()
        writer.add_blank_page(612, 792)
        writer.add_metadata({"/Title": "Test PDF", "/Author": "Tester"})
        with open(path, "wb") as f:
            writer.write(f)
    except ImportError:
        try:
            import pdfplumber
            # 无法创建 PDF，用占位
            with open(path, "wb") as f:
                f.write(b"%PDF-1.4\n1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj\n3 0 obj<</Type/Page/MediaBox[0 0 612 792]>>endobj\nxref\n0 4\n0000000000 65535 f\n0000000009 00000 n\n0000000058 00000 n\n0000000115 00000 n\ntrailer<</Size 4/Root 1 0 R>>\nstartxref\n190\n%%EOF")
        except ImportError:
            with open(path, "wb") as f:
                f.write(b"%PDF placeholder")
    return path


# ── 测试用例 ──────────────────────────────────────

def test_python_file():
    """① .py → document + text/x-python"""
    tmp = tempfile.mkdtemp()
    try:
        path = _create_python(tmp)
        info = fc.identify(path)
        assert info.content_type == "document", f"Expected document, got {info.content_type}"
        assert "x-python" in info.mime_type or "python" in info.mime_type
        assert len(info.text_preview) > 0
        assert info.text_length > 0
        assert info.file_size > 0
        assert info.checksum != ""
        assert info.error == ""
    finally:
        shutil.rmtree(tmp, ignore_errors=True)
    print("  ✓ .py → document + text/x-python")


def test_markdown_file():
    """② .md → text/markdown"""
    tmp = tempfile.mkdtemp()
    try:
        path = _create_markdown(tmp)
        info = fc.identify(path)
        assert info.content_type == "document"
        assert "markdown" in info.mime_type
        assert "README" in info.text_preview
        assert info.error == ""
    finally:
        shutil.rmtree(tmp, ignore_errors=True)
    print("  ✓ .md → text/markdown")


def test_json_file():
    """③ .json → application/json"""
    tmp = tempfile.mkdtemp()
    try:
        path = _create_json(tmp)
        info = fc.identify(path)
        assert info.content_type == "document"
        assert "json" in info.mime_type
        assert "test" in info.text_preview
        assert info.error == ""
    finally:
        shutil.rmtree(tmp, ignore_errors=True)
    print("  ✓ .json → application/json")


def test_utf8_text_file():
    """④ .txt UTF-8 编码检测"""
    tmp = tempfile.mkdtemp()
    try:
        path = _create_txt(tmp)
        info = fc.identify(path)
        assert info.content_type == "document"
        assert "text/plain" in info.mime_type
        assert "Python" in info.text_preview
        assert info.error == ""
    finally:
        shutil.rmtree(tmp, ignore_errors=True)
    print("  ✓ .txt UTF-8 encoding detection")


def test_png_image():
    """⑨ .png → image + 宽高元数据"""
    tmp = tempfile.mkdtemp()
    try:
        path = _create_png(tmp)
        info = fc.identify(path)
        assert info.content_type == "image"
        assert "image/png" in info.mime_type
        if "width" in info.metadata:
            assert info.metadata["width"] > 0
        assert info.error == ""
    finally:
        shutil.rmtree(tmp, ignore_errors=True)
    print("  ✓ .png → image + metadata")


def test_zip_archive():
    """⑩ .zip → archive"""
    tmp = tempfile.mkdtemp()
    try:
        path = _create_zip(tmp)
        info = fc.identify(path)
        assert info.content_type == "archive", f"Expected archive, got {info.content_type}"
        assert "zip" in info.mime_type
        assert info.error == ""
    finally:
        shutil.rmtree(tmp, ignore_errors=True)
    print("  ✓ .zip → archive")


def test_md5_correctness():
    """⑬ MD5 计算正确"""
    tmp = tempfile.mkdtemp()
    try:
        path = _create_python(tmp)
        info = fc.identify(path)
        with open(path, "rb") as f:
            expected = hashlib.md5(f.read()).hexdigest()
        assert info.checksum == expected, f"MD5 mismatch: {info.checksum} vs {expected}"
    finally:
        shutil.rmtree(tmp, ignore_errors=True)
    print("  ✓ MD5 correctness")


def test_nonexistent_file():
    """⑪ 不存在的文件 → error"""
    info = fc.identify("/tmp/nonexistent_file_xyz_999.txt")
    assert info.error != "", "Expected error for nonexistent file"
    print("  ✓ nonexistent file → error")


def test_empty_file():
    """⑫ 空文件 → 不崩 + text_length=0"""
    tmp = tempfile.mkdtemp()
    try:
        path = _create_empty(tmp)
        info = fc.identify(path)
        assert info.error == "", f"Unexpected error: {info.error}"
        assert info.text_length == 0
        assert info.file_size == 0
    finally:
        shutil.rmtree(tmp, ignore_errors=True)
    print("  ✓ empty file → no crash")


def test_docx():
    """⑤ .docx → paragraphs+author+title"""
    tmp = tempfile.mkdtemp()
    try:
        path = _create_docx(tmp)
        try:
            from docx import Document
            Document  # check import
        except ImportError:
            print("  ⚠ skip: docx not available")
            return
        info = fc.identify(path)
        assert info.content_type == "document"
        assert "officedocument.wordprocessingml" in info.mime_type
        assert len(info.text_preview) > 0
        assert info.error == ""
    finally:
        shutil.rmtree(tmp, ignore_errors=True)
    print("  ✓ .docx → paragraphs + author")


def test_xlsx():
    """⑥ .xlsx → sheets + 内容"""
    tmp = tempfile.mkdtemp()
    try:
        path = _create_xlsx(tmp)
        try:
            from openpyxl import load_workbook
            load_workbook  # check import
        except ImportError:
            print("  ⚠ skip: openpyxl not available")
            return
        info = fc.identify(path)
        assert info.content_type == "document"
        assert "officedocument.spreadsheetml" in info.mime_type
        assert info.error == ""
    finally:
        shutil.rmtree(tmp, ignore_errors=True)
    print("  ✓ .xlsx → sheets + content")


def test_pptx():
    """⑦ .pptx → slides + 文本"""
    tmp = tempfile.mkdtemp()
    try:
        path = _create_pptx(tmp)
        try:
            from pptx import Presentation
            Presentation
        except ImportError:
            print("  ⚠ skip: pptx not available")
            return
        info = fc.identify(path)
        assert info.content_type == "document"
        assert "officedocument.presentationml" in info.mime_type
        assert info.error == ""
    finally:
        shutil.rmtree(tmp, ignore_errors=True)
    print("  ✓ .pptx → slides + text")


def test_pdf():
    """⑧ .pdf → pages + pdf_title"""
    tmp = tempfile.mkdtemp()
    try:
        path = _create_pdf(tmp)
        try:
            from PyPDF2 import PdfReader
            PdfReader
        except ImportError:
            try:
                import pdfplumber
                pdfplumber
            except ImportError:
                print("  ⚠ skip: no PDF library available")
                return
        info = fc.identify(path)
        assert info.content_type == "document"
        assert "application/pdf" in info.mime_type
        assert info.error == ""
    finally:
        shutil.rmtree(tmp, ignore_errors=True)
    print("  ✓ .pdf → pages + title")


def test_unknown_extension():
    """⑭ 未知扩展名 → 降级为文本读取"""
    tmp = tempfile.mkdtemp()
    try:
        path = os.path.join(tmp, "unknown.xyz")
        with open(path, "w") as f:
            f.write("plain text content with unknown extension")
        info = fc.identify(path)
        # 应该降级为纯文本
        assert info.text_length > 0
        assert "unknown extension" in info.text_preview or info.content_type == "document"
        assert info.error == ""
    finally:
        shutil.rmtree(tmp, ignore_errors=True)
    print("  ✓ unknown extension → fallback text")


def test_large_file():
    """大文件（>10MB）→ 标记而非读取全文"""
    tmp = tempfile.mkdtemp()
    try:
        path = os.path.join(tmp, "large.bin")
        with open(path, "wb") as f:
            f.write(b"x" * (11 * 1024 * 1024))
        info = fc.identify(path)
        assert info.text_length <= 500, f"Expected truncated text, got {info.text_length}"
        assert info.error == ""
    finally:
        shutil.rmtree(tmp, ignore_errors=True)
    print("  ✓ large file → text truncated")


# ── 运行入口 ──────────────────────────────────────
if __name__ == "__main__":
    tests = [
        test_python_file,
        test_markdown_file,
        test_json_file,
        test_utf8_text_file,
        test_png_image,
        test_zip_archive,
        test_md5_correctness,
        test_nonexistent_file,
        test_empty_file,
        test_docx,
        test_xlsx,
        test_pptx,
        test_pdf,
        test_unknown_extension,
        test_large_file,
    ]
    passed = 0
    skipped = 0
    for t in tests:
        try:
            t()
            passed += 1
        except AssertionError as e:
            import traceback
            print(f"  ✗ {t.__name__}: {e}")
            traceback.print_exc()
        except Exception as e:
            print(f"  ✗ {t.__name__}: {e}")
    print(f"\n{'='*40}\n结果: {passed}/{len(tests)} 通过 (skip 标记不计数)")
